"""
Weekly Summary Generator for VitaInspire AI Research

Generates a weekly summary from daily reports every Sunday with:
- HTML summary report
- Read aloud audio version (using gTTS)
- Slide deck (PowerPoint)
- Infographic (PNG image)

All outputs are saved to Google Drive.
"""

import os
import re
import datetime
from pathlib import Path
import sys
sys.path.insert(0, str(Path(__file__).parent))

try:
    import google.generativeai as genai
except ImportError:
    print("Please install: pip install google-generativeai")
    
from google_drive_config import (
    get_google_drive_service,
    setup_vitainspire_folders,
    upload_content_to_drive,
    upload_file_to_drive
)

SCRIPT_DIR = Path(__file__).parent
BLOG_FILE_PATH = SCRIPT_DIR.parent / "js" / "blog-posts.js"
TEMP_DIR = SCRIPT_DIR / "temp_weekly"
API_KEY = os.environ.get("GEMINI_API_KEY")


def parse_blog_posts():
    """Parse blog posts from blog-posts.js."""
    with open(BLOG_FILE_PATH, 'r', encoding='utf-8') as f:
        content = f.read()
    
    posts = []
    pattern = re.compile(
        r'\{\s*id:\s*"([^"]+)",\s*title:\s*"([^"]+)",\s*'
        r'date:\s*"([^"]+)",\s*category:\s*"([^"]+)",\s*'
        r'author:\s*"([^"]+)",\s*excerpt:\s*"([^"]+)",\s*'
        r'content:\s*`([^`]+)`', re.DOTALL
    )
    
    for m in pattern.finditer(content):
        posts.append({
            'id': m.group(1), 'title': m.group(2), 'date': m.group(3),
            'category': m.group(4), 'author': m.group(5),
            'excerpt': m.group(6), 'content': m.group(7).strip()
        })
    return posts


def get_week_posts(posts):
    """Get posts from the past 7 days."""
    today = datetime.datetime.now()
    week_ago = today - datetime.timedelta(days=7)
    
    week_posts = []
    for post in posts:
        try:
            post_date = datetime.datetime.strptime(post['date'], "%B %d, %Y")
            if week_ago <= post_date <= today:
                week_posts.append(post)
        except ValueError:
            continue
    return week_posts


def generate_weekly_summary(posts):
    """Use Gemini to generate a weekly summary."""
    if not API_KEY:
        print("Warning: GEMINI_API_KEY not set, using simple summary")
        return create_simple_summary(posts)
    
    genai.configure(api_key=API_KEY)
    model = genai.GenerativeModel('gemini-2.5-flash')
    
    # Prepare post summaries
    post_summaries = "\n\n".join([
        f"Title: {p['title']}\nCategory: {p['category']}\nExcerpt: {p['excerpt']}"
        for p in posts[:20]  # Limit to avoid token limits
    ])
    
    prompt = f"""You are an expert AI research analyst for VitaInspire, a social impact organization.
Create a comprehensive weekly summary from these daily AI research articles about AI for social good in India.

Articles from this week:
{post_summaries}

Generate a JSON response with these fields:
- headline: An engaging headline for the weekly summary
- executive_summary: 2-3 paragraph overview of key themes and takeaways
- key_themes: Array of 3-5 major themes with title and description
- top_stories: Array of 3 most impactful stories with title and why_it_matters
- stats: Object with total_articles, categories_covered (array), sectors_impacted
- outlook: 1 paragraph on what to watch next week
- audio_script: A 2-minute read-aloud friendly version of the summary
- slide_bullets: Array of 5-6 key points for a presentation"""

    try:
        response = model.generate_content(prompt, generation_config={
            "response_mime_type": "application/json"
        })
        import json
        return json.loads(response.text)
    except Exception as e:
        print(f"AI summary generation failed: {e}")
        return create_simple_summary(posts)


def create_simple_summary(posts):
    """Create a comprehensive summary without AI - includes full article details."""
    import re
    
    categories = list(set(p['category'] for p in posts))
    
    # Strip HTML tags for text content
    def strip_html(text):
        clean = re.sub(r'<[^>]+>', ' ', text)
        clean = re.sub(r'\s+', ' ', clean).strip()
        return clean
    
    # Build comprehensive executive summary
    exec_parts = [f"This week, VitaInspire published {len(posts)} AI research articles covering {', '.join(categories[:5])}."]
    for i, post in enumerate(posts[:5], 1):
        exec_parts.append(f"\n\n{i}. **{post['title']}** ({post['category']}): {post['excerpt']}")
    executive_summary = ''.join(exec_parts)
    
    # Build comprehensive audio script (2-3 minute read)
    audio_parts = [
        f"Welcome to VitaInspire's Weekly AI Research Summary. ",
        f"This week, we covered {len(posts)} articles exploring AI for social good in India. ",
        f"The topics included {', '.join(categories[:4])}. ",
        "Let me walk you through the highlights. "
    ]
    
    for i, post in enumerate(posts, 1):
        content_clean = strip_html(post['content'])[:400]  # First 400 chars of content
        audio_parts.append(
            f"\n\nArticle {i}: {post['title']}. "
            f"In the {post['category']} sector: {post['excerpt']} "
            f"{content_clean}... "
        )
    
    audio_parts.append(
        "\n\nThat concludes this week's AI research roundup from VitaInspire. "
        "Join us next week for more insights on AI transforming social impact in India. "
        "Building AI careers, transforming social impact."
    )
    audio_script = ''.join(audio_parts)
    
    # Build detailed slide bullets with content
    slide_bullets = []
    for post in posts:
        content_clean = strip_html(post['content'])[:200]
        slide_bullets.append({
            'title': post['title'],
            'category': post['category'],
            'excerpt': post['excerpt'],
            'content_preview': content_clean
        })
    
    # Build key themes with actual descriptions from posts
    theme_dict = {}
    for post in posts:
        cat = post['category']
        if cat not in theme_dict:
            theme_dict[cat] = {'title': cat, 'description': post['excerpt'], 'articles': []}
        theme_dict[cat]['articles'].append(post['title'])
    
    key_themes = []
    for cat, data in theme_dict.items():
        desc = f"{data['description']} (Articles: {', '.join(data['articles'][:2])})"
        key_themes.append({'title': cat, 'description': desc})
    
    # Build top stories with full excerpts
    top_stories = []
    for post in posts[:5]:
        content_clean = strip_html(post['content'])[:300]
        top_stories.append({
            'title': post['title'],
            'why_it_matters': f"{post['excerpt']} {content_clean}...",
            'category': post['category'],
            'date': post['date']
        })
    
    return {
        "headline": f"Weekly AI Research Roundup: {len(posts)} Articles on AI for Social Good",
        "executive_summary": executive_summary,
        "key_themes": key_themes[:6],
        "top_stories": top_stories,
        "stats": {"total_articles": len(posts), "categories_covered": categories},
        "outlook": f"Next week, continue monitoring developments in {', '.join(categories[:3])} as AI continues to transform social impact across India. Stay tuned for more research on how AI is empowering communities and creating positive change.",
        "audio_script": audio_script,
        "slide_bullets": slide_bullets,
        "full_posts": posts  # Include full posts for detailed outputs
    }


def create_weekly_html(summary, week_start, week_end):
    """Generate HTML report for weekly summary."""
    themes_html = "".join([
        f'<div class="theme"><h4>{t["title"]}</h4><p>{t.get("description","")}</p></div>'
        for t in summary.get("key_themes", [])
    ])
    
    stories_html = "".join([
        f'<div class="story"><h4>{s["title"]}</h4><p>{s.get("why_it_matters","")}</p></div>'
        for s in summary.get("top_stories", [])
    ])
    
    return f'''<!DOCTYPE html>
<html><head><meta charset="UTF-8"><title>VitaInspire Weekly Summary - {week_end}</title>
<style>
body{{font-family:'Segoe UI',sans-serif;background:linear-gradient(135deg,#0f0f23,#1a1a3e);color:#e0e0e0;padding:40px;}}
.container{{max-width:1000px;margin:auto;}}
.header{{text-align:center;background:linear-gradient(135deg,#667eea,#764ba2);padding:50px;border-radius:25px;margin-bottom:40px;}}
h1{{font-size:2.5rem;margin-bottom:10px;}}
.period{{font-size:1.1rem;opacity:0.9;}}
.section{{background:#1a1a3e;padding:35px;border-radius:20px;margin-bottom:30px;border:1px solid rgba(255,255,255,0.1);}}
.section h3{{color:#00d4aa;margin-bottom:20px;font-size:1.4rem;}}
.summary{{font-size:1.1rem;line-height:1.8;}}
.theme,.story{{background:rgba(0,0,0,0.3);padding:20px;border-radius:12px;margin-bottom:15px;}}
.theme h4,.story h4{{color:#fff;margin-bottom:10px;}}
.stats{{display:flex;gap:20px;flex-wrap:wrap;}}
.stat{{flex:1;min-width:150px;background:linear-gradient(135deg,#667eea,#764ba2);padding:25px;border-radius:15px;text-align:center;}}
.stat-number{{font-size:2rem;font-weight:bold;}}
.stat-label{{opacity:0.9;}}
.footer{{text-align:center;padding:30px;color:#888;}}
</style></head>
<body><div class="container">
<div class="header"><h1>📊 Weekly AI Research Summary</h1>
<p class="period">{week_start} - {week_end}</p></div>

<div class="section"><h3>Executive Summary</h3>
<div class="summary">{summary.get("executive_summary","")}</div></div>

<div class="section"><h3>📈 This Week's Stats</h3>
<div class="stats">
<div class="stat"><div class="stat-number">{summary.get("stats",{}).get("total_articles",0)}</div>
<div class="stat-label">Articles</div></div>
<div class="stat"><div class="stat-number">{len(summary.get("stats",{}).get("categories_covered",[]))}</div>
<div class="stat-label">Categories</div></div></div></div>

<div class="section"><h3>🎯 Key Themes</h3>{themes_html}</div>
<div class="section"><h3>⭐ Top Stories</h3>{stories_html}</div>
<div class="section"><h3>🔮 Week Ahead</h3><p>{summary.get("outlook","")}</p></div>

<footer class="footer"><p><strong>VitaAInspire</strong> - Building AI Careers, Transforming Social Impact</p>
<p>© 2026 VitaInspire Private Limited</p></footer>
</div></body></html>'''


def create_audio_file(summary, filename):
    """Generate audio file using gTTS."""
    try:
        from gtts import gTTS
        script = summary.get("audio_script", summary.get("executive_summary", ""))
        tts = gTTS(text=script, lang='en', slow=False)
        tts.save(filename)
        print(f"✅ Audio file created: {filename}")
        return True
    except ImportError:
        print("gTTS not installed. Run: pip install gtts")
        return False
    except Exception as e:
        print(f"Audio generation failed: {e}")
        return False


def create_slide_deck(summary, filename):
    """Generate PowerPoint slide deck with comprehensive article content."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as e:
        print(f"python-pptx not installed. Run: pip install python-pptx. Error: {e}")
        return False
    
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
        tf = title.text_frame
        tf.text = summary.get("headline", "Weekly AI Research Summary")
        tf.paragraphs[0].font.size = Pt(40)
        tf.paragraphs[0].font.bold = True
        
        subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
        subtitle.text_frame.text = "VitaInspire - Building AI Careers, Transforming Social Impact"
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        
        stats_text = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(1))
        stats = summary.get("stats", {})
        stats_text.text_frame.text = f"{stats.get('total_articles', 0)} Articles | {len(stats.get('categories_covered', []))} Categories"
        stats_text.text_frame.paragraphs[0].font.size = Pt(20)
        
        # Executive Summary slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title.text_frame.text = "📊 Executive Summary"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        
        exec_text = summary.get("executive_summary", "")[:800]
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5))
        body.text_frame.text = exec_text
        body.text_frame.paragraphs[0].font.size = Pt(16)
        body.text_frame.word_wrap = True
        
        # Individual article slides
        bullets = summary.get("slide_bullets", [])
        for i, article in enumerate(bullets, 1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Handle both dict format (new) and string format (old)
            if isinstance(article, dict):
                article_title = article.get('title', 'Article')
                category = article.get('category', '')
                excerpt = article.get('excerpt', '')
                content = article.get('content_preview', '')
            else:
                article_title = str(article)
                category = ''
                excerpt = ''
                content = ''
            
            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            title_box.text_frame.text = f"Article {i}: {article_title[:60]}..."
            title_box.text_frame.paragraphs[0].font.size = Pt(26)
            title_box.text_frame.paragraphs[0].font.bold = True
            
            # Category badge
            if category:
                cat_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(0.5))
                cat_box.text_frame.text = f"🏷️ {category}"
                cat_box.text_frame.paragraphs[0].font.size = Pt(18)
            
            # Excerpt
            if excerpt:
                excerpt_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.5))
                excerpt_box.text_frame.text = f"Summary: {excerpt[:250]}"
                excerpt_box.text_frame.paragraphs[0].font.size = Pt(18)
                excerpt_box.text_frame.word_wrap = True
            
            # Content preview
            if content:
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.3), Inches(3.5))
                content_box.text_frame.text = f"Key Points: {content[:400]}..."
                content_box.text_frame.paragraphs[0].font.size = Pt(16)
                content_box.text_frame.word_wrap = True
        
        # Outlook slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        title.text_frame.text = "🔮 Looking Ahead"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        
        outlook_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(4))
        outlook_box.text_frame.text = summary.get("outlook", "")
        outlook_box.text_frame.paragraphs[0].font.size = Pt(20)
        outlook_box.text_frame.word_wrap = True
        
        # Closing slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        closing = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        closing.text_frame.text = "Thank You!\n\nVitaInspire\nBuilding AI Careers | Transforming Social Impact"
        closing.text_frame.paragraphs[0].font.size = Pt(36)
        closing.text_frame.paragraphs[0].font.bold = True
        
        prs.save(filename)
        print(f"✅ Slide deck created: {filename}")
        return True
    except Exception as e:
        print(f"Slide deck generation failed: {e}")
        import traceback
        traceback.print_exc()
        return False


def create_infographic(summary, filename):
    """Generate comprehensive infographic image with full article details."""
    try:
        from PIL import Image, ImageDraw, ImageFont
        import textwrap
        
        # Calculate height based on number of articles
        num_articles = len(summary.get("top_stories", []))
        base_height = 800
        per_article_height = 250
        height = base_height + (num_articles * per_article_height)
        width = 1200
        
        img = Image.new('RGB', (width, height), color='#1a1a2e')
        draw = ImageDraw.Draw(img)
        
        try:
            font_title = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 48)
            font_heading = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 32)
            font_subhead = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 26)
            font_body = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
            font_small = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 16)
        except:
            font_title = font_heading = font_subhead = font_body = font_small = ImageFont.load_default()
        
        # Gradient header
        for i in range(220):
            r = int(102 + (118 - 102) * i / 220)
            g = int(126 + (75 - 126) * i / 220)
            b = int(234 + (162 - 234) * i / 220)
            draw.line([(0, i), (width, i)], fill=(r, g, b))
        
        # Title
        draw.text((50, 40), "VitaInspire Weekly", fill='white', font=font_title)
        draw.text((50, 100), "AI Research Summary", fill='white', font=font_title)
        
        # Stats bar
        stats = summary.get("stats", {})
        stats_text = f"📊 {stats.get('total_articles', 0)} Articles  |  🏷️ {len(stats.get('categories_covered', []))} Categories"
        draw.text((50, 175), stats_text, fill='#e0e0e0', font=font_subhead)
        
        y = 260
        
        # Executive Summary section
        draw.rounded_rectangle([(30, y), (width-30, y+120)], radius=15, fill='#252545')
        draw.text((50, y+15), "📋 This Week's Highlights", fill='#00d4aa', font=font_heading)
        exec_text = summary.get("executive_summary", "")[:200] + "..."
        wrapped = textwrap.wrap(exec_text, width=70)
        text_y = y + 55
        for line in wrapped[:2]:
            draw.text((50, text_y), line, fill='#e0e0e0', font=font_body)
            text_y += 28
        
        y += 150
        
        # Articles section
        draw.text((50, y), "📰 Featured Articles", fill='#667eea', font=font_heading)
        y += 50
        
        for i, story in enumerate(summary.get("top_stories", []), 1):
            # Article card background
            draw.rounded_rectangle([(30, y), (width-30, y+200)], radius=15, fill='#252545')
            
            # Article number
            draw.ellipse([(50, y+15), (90, y+55)], fill='#667eea')
            draw.text((63, y+22), str(i), fill='white', font=font_subhead)
            
            # Title
            title = story.get('title', '')[:70]
            draw.text((110, y+20), title, fill='white', font=font_subhead)
            
            # Category
            category = story.get('category', '')
            if category:
                draw.text((110, y+55), f"🏷️ {category}", fill='#00d4aa', font=font_small)
            
            # Why it matters (excerpt + content preview)
            why = story.get('why_it_matters', '')[:280]
            wrapped = textwrap.wrap(why, width=75)
            text_y = y + 85
            for line in wrapped[:4]:
                draw.text((50, text_y), line, fill='#b0b0b0', font=font_small)
                text_y += 24
            
            y += 220
        
        # Outlook section
        y += 20
        draw.rounded_rectangle([(30, y), (width-30, y+100)], radius=15, fill='#252545')
        draw.text((50, y+15), "🔮 Looking Ahead", fill='#667eea', font=font_heading)
        outlook = summary.get("outlook", "")[:150]
        wrapped = textwrap.wrap(outlook, width=70)
        text_y = y + 55
        for line in wrapped[:2]:
            draw.text((50, text_y), line, fill='#e0e0e0', font=font_body)
            text_y += 28
        
        # Footer
        y = height - 80
        draw.rectangle([(0, y), (width, height)], fill='#0f0f23')
        draw.text((50, y+20), "VitaInspire", fill='#667eea', font=font_heading)
        draw.text((250, y+25), "Building AI Careers  •  Transforming Social Impact", fill='#888888', font=font_body)
        
        img.save(filename)
        print(f"✅ Infographic created: {filename}")
        return True
    except ImportError:
        print("Pillow not installed. Run: pip install Pillow")
        return False
    except Exception as e:
        print(f"Infographic generation failed: {e}")
        import traceback
        traceback.print_exc()
        return False


def generate_weekly_summary_and_upload():
    """Main function to generate all weekly summary formats and upload."""
    today = datetime.datetime.now()
    
    # Check if today is Sunday
    if today.weekday() != 6:
        print(f"Today is {today.strftime('%A')}, not Sunday. Running anyway...")
    
    # Calculate week range
    week_end = today.strftime("%B %d, %Y")
    week_start = (today - datetime.timedelta(days=6)).strftime("%B %d, %Y")
    
    print(f"\n📅 Generating weekly summary for {week_start} to {week_end}\n")
    
    # Get week's posts
    all_posts = parse_blog_posts()
    week_posts = get_week_posts(all_posts)
    
    if not week_posts:
        print("No posts found for this week!")
        return
    
    print(f"Found {len(week_posts)} posts for this week")
    
    # Generate summary
    print("\n🤖 Generating AI-powered summary...")
    summary = generate_weekly_summary(week_posts)
    
    # Create temp directory
    TEMP_DIR.mkdir(exist_ok=True)
    
    # Connect to Drive
    print("\n☁️ Connecting to Google Drive...")
    service = get_google_drive_service()
    folders = setup_vitainspire_folders(service)
    
    date_suffix = today.strftime("%Y%m%d")
    
    # 1. HTML Report
    print("\n📄 Creating HTML report...")
    html_content = create_weekly_html(summary, week_start, week_end)
    upload_content_to_drive(
        service, html_content,
        f"Weekly_Summary_{date_suffix}.html",
        folders['weekly_html'], 'text/html'
    )
    
    # 2. Audio (Read Aloud)
    print("\n🎧 Creating audio version...")
    audio_file = TEMP_DIR / f"Weekly_Summary_{date_suffix}.mp3"
    if create_audio_file(summary, str(audio_file)):
        upload_file_to_drive(service, audio_file, folders['weekly_audio'])
    
    # 3. Slide Deck
    print("\n📊 Creating slide deck...")
    slides_file = TEMP_DIR / f"Weekly_Summary_{date_suffix}.pptx"
    if create_slide_deck(summary, str(slides_file)):
        upload_file_to_drive(service, slides_file, folders['weekly_slides'])
    
    # 4. Infographic
    print("\n🎨 Creating infographic...")
    infographic_file = TEMP_DIR / f"Weekly_Summary_{date_suffix}.png"
    if create_infographic(summary, str(infographic_file)):
        upload_file_to_drive(service, infographic_file, folders['weekly_infographics'])
    
    # Cleanup
    import shutil
    if TEMP_DIR.exists():
        shutil.rmtree(TEMP_DIR)
    
    print("\n" + "="*50)
    print("✅ Weekly summary generation complete!")
    print("="*50)


if __name__ == "__main__":
    generate_weekly_summary_and_upload()
